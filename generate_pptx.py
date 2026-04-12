"""Generate DistrictPilot AI Hackathon PPT (final deck, dark navy theme)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
NAVY = RGBColor(0x1E, 0x27, 0x61)
DARK_NAVY = RGBColor(0x14, 0x1B, 0x41)
ICE_BLUE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x29, 0xB6, 0xF6)
LIGHT_GRAY = RGBColor(0xA0, 0xA8, 0xBE)
ORANGE = RGBColor(0xFF, 0x9F, 0x43)
GREEN = RGBColor(0x66, 0xBB, 0x6A)

BLANK = prs.slide_layouts[6]  # blank layout

def add_bg(slide, color=NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, text, x, y, w, h, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tf

def add_multiline(slide, lines, x, y, w, h, size=16, color=WHITE, bold=False, spacing=1.2, bullet=False):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Calibri"
        p.space_after = Pt(size * spacing * 0.4)
    return tf

def add_rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))  # MSO_SHAPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_table_slide(slide, headers, rows, x, y, w, row_h=0.45):
    cols = len(headers)
    total_rows = len(rows) + 1
    table_shape = slide.shapes.add_table(total_rows, cols, Inches(x), Inches(y), Inches(w), Inches(row_h * total_rows))
    table = table_shape.table
    col_w = w / cols
    for i in range(cols):
        table.columns[i].width = Inches(col_w)
    # Headers
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x29, 0x4D, 0x8F)
    # Data
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = WHITE
                p.font.name = "Calibri"
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x1A, 0x23, 0x50) if i % 2 == 0 else RGBColor(0x22, 0x2D, 0x5E)
    return table

# ═══════════════════════════════════════
# Slide 1: Title
# ═══════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, DARK_NAVY)
add_rect(s, 0, 0, 13.333, 0.08, ACCENT)
add_text(s, "DistrictPilot AI", 0.8, 1.8, 11.5, 1.5, size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "이사 직후 72시간 골든타임을 예측하는 Snowflake-Native AI 에이전트", 0.8, 3.2, 11.5, 0.8, size=24, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s, "서초 · 영등포 · 중구  홈서비스 수요 예측 → 구별 집행 액션 추천", 0.8, 4.1, 11.5, 0.6, size=18, color=ICE_BLUE, align=PP_ALIGN.CENTER)
add_text(s, "Snowflake Korea Hackathon 2026  |  Tech Track", 0.8, 5.2, 11.5, 0.5, size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(s, "Doeon Kim  |  github.com/KIM3310", 0.8, 6.2, 11.5, 0.4, size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# Slide 2: Problem
# ═══════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
add_rect(s, 0, 0, 0.08, 7.5, ACCENT)
add_text(s, "Problem", 0.5, 0.4, 5, 0.6, size=36, bold=True, color=ACCENT)
add_text(s, '"다음 달 서초·영등포·중구 중\n어디에서 전입 직후 수요를 먼저 잡을 것인가?"', 0.8, 1.5, 8, 1.5, size=24, bold=True, color=WHITE)
add_multiline(s, [
    "기존: 경험과 감에 의존한 지역 집행",
    "문제: 전입·이사 신호와 소비 반응, 운영 제약을 같이 읽기 어려움",
    "결과: 타이밍 미스, 과집행/미집행, 설치·운영 리스크 누적",
], 0.8, 3.5, 7, 2, size=17, color=ICE_BLUE)
add_rect(s, 9, 1.5, 3.5, 4.5, RGBColor(0x22, 0x2D, 0x5E))
add_text(s, "DistrictPilot AI", 9.2, 1.8, 3.1, 0.5, size=18, bold=True, color=ACCENT)
add_multiline(s, [
    "전입 신호 감지",
    "수요 근거 제시",
    "실시간 갱신",
    "운영 액션 연결",
], 9.3, 2.5, 2.9, 3, size=15, color=WHITE)

# ═══════════════════════════════════════
# Slide 3: Data Sources
# ═══════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
add_rect(s, 0, 0, 0.08, 7.5, ACCENT)
add_text(s, "Data Sources", 0.5, 0.4, 5, 0.6, size=36, bold=True, color=ACCENT)
add_text(s, "Sponsor (Snowflake Marketplace)", 0.5, 1.2, 6, 0.5, size=20, bold=True, color=WHITE)
add_table_slide(s, ["Source", "Data", "Level"], [
    ["SPH", "유동인구 / 카드소비 / 자산소득", "법정동 → 구"],
    ["Richgo", "매매/전세 시세 / 인구이동", "시군구"],
    ["AJD (Optional)", "렌탈 / 통신 / 마케팅 / CS", "별도 통합 시 사용"],
], 0.5, 1.8, 6)
add_text(s, "전입 신호 + 소비 반응 + 운영 액션을 한 문제로 묶는 구조", 7.0, 1.8, 5.4, 0.6, size=18, color=ORANGE)
add_text(s, "External (Public Open Data)", 0.5, 3.8, 6, 0.5, size=20, bold=True, color=WHITE)
add_table_slide(s, ["Source", "Data", "License"], [
    ["한국천문연구원", "공휴일/특일 캘린더", "공공누리 1유형"],
    ["행정안전부", "연령·성별 주민등록 인구", "공공누리 1유형"],
    ["한국관광 데이터랩", "관광수요/외래객 지수", "공공누리 1유형"],
    ["서울시 상권분석", "상권변화지표", "공공누리 1유형"],
], 0.5, 4.4, 6)

# ═══════════════════════════════════════
# Slide 4: Architecture
# ═══════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
add_rect(s, 0, 0, 0.08, 7.5, ACCENT)
add_text(s, "Architecture", 0.5, 0.4, 5, 0.6, size=36, bold=True, color=ACCENT)
# Architecture layers as boxes
layers = [
    ("Data Sources", "SPH + Richgo + 공휴일 + 연령 + 관광 + 상권 (AJD optional)", 0.8, 1.3),
    ("Feature Store", "DT_FEATURE_MART -> FEATURE_MART_V2", 0.8, 2.3),
    ("ML / AI", "DISTRICTPILOT_FORECAST_V2 + Semantic View + Cortex Search + AI_COMPLETE", 0.8, 3.3),
    ("Decision", "AI_COMPLETE Structured Output -> Capture Plan + Operator Actions", 0.8, 4.3),
    ("App", "Streamlit in Snowflake (Capture Plan | Signals | AI | Scenario | Ops)", 0.8, 5.3),
    ("Ops", "Tasks (daily refresh / weekly retrain) + V_APP_HEALTH + Query Tags", 0.8, 6.3),
]
for label, desc, x, y in layers:
    add_rect(s, x, y, 11.5, 0.8, RGBColor(0x22, 0x2D, 0x5E))
    add_rect(s, x, y, 0.06, 0.8, ACCENT)
    add_text(s, label, x + 0.2, y + 0.05, 2.5, 0.35, size=14, bold=True, color=ACCENT)
    add_text(s, desc, x + 2.8, y + 0.05, 8.5, 0.7, size=13, color=ICE_BLUE)

# ═══════════════════════════════════════
# Slide 5: Ablation Study
# ═══════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
add_rect(s, 0, 0, 0.08, 7.5, ACCENT)
add_text(s, "Ablation Study", 0.5, 0.4, 5, 0.6, size=36, bold=True, color=ACCENT)
add_text(s, "외부 데이터 추가로 정량적 MAPE 개선 증명", 0.5, 1.0, 8, 0.5, size=18, color=ICE_BLUE)
add_table_slide(s, ["Model", "Features", "Description"], [
    ["A", "Y-only", "스폰서 매출 데이터만 (baseline)"],
    ["B", "+ Holiday", "공휴일/연휴/영업일수 추가"],
    ["C", "+ Demographics", "20-39세/시니어 연령구조 추가"],
    ["D", "+ Tourism", "관광수요/외래객 지수 추가"],
    ["E", "Full", "상권안정도/점포변동 추가 (production)"],
], 0.5, 1.8, 7)
add_rect(s, 8, 1.8, 4.5, 4, RGBColor(0x22, 0x2D, 0x5E))
add_text(s, "MAPE Improvement", 8.2, 2.0, 4, 0.5, size=16, bold=True, color=ACCENT)
add_multiline(s, [
    "Model A → E:",
    "MAPE 단계적 감소",
    "",
    "SHOW_EVALUATION_METRICS()",
    "로 cross-validation 기반",
    "out-of-sample 평가",
], 8.3, 2.6, 3.8, 2.5, size=14, color=ICE_BLUE)
add_text(s, "앱 Allocation 탭에서 실시간 확인", 8.3, 5.0, 3.8, 0.5, size=13, color=ORANGE)

# ═══════════════════════════════════════
# Slide 6: Feature Importance
# ═══════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
add_rect(s, 0, 0, 0.08, 7.5, ACCENT)
add_text(s, "Feature Importance", 0.5, 0.4, 6, 0.6, size=36, bold=True, color=ACCENT)
add_text(s, "EXPLAIN_FEATURE_IMPORTANCE() — 왜 이 구가 높은가?", 0.5, 1.0, 10, 0.5, size=18, color=ICE_BLUE)
# 3 district columns
for i, (district, features) in enumerate([
    ("서초구", ["TOTAL_SALES lag", "HOLIDAY_DAYS", "AGE_20_39_SHARE", "STABILITY_SCORE"]),
    ("영등포구", ["NET_MOVE", "WORK_POP", "TOURISM_IDX", "COFFEE sales"]),
    ("중구", ["FOREIGN_VISITOR_IDX", "VISIT_POP", "ACCOMMODATION", "NET_STORE_CHANGE"]),
]):
    x = 0.5 + i * 4.2
    add_rect(s, x, 1.8, 3.8, 4.5, RGBColor(0x22, 0x2D, 0x5E))
    add_rect(s, x, 1.8, 3.8, 0.06, ACCENT)
    add_text(s, district, x + 0.2, 2.0, 3.4, 0.5, size=20, bold=True, color=WHITE)
    add_multiline(s, [f"{j+1}. {f}" for j, f in enumerate(features)], x + 0.3, 2.7, 3.2, 3, size=15, color=ICE_BLUE)
add_text(s, "Feature importance 0~1 범위 정규화, 상대적 중요도 기반", 0.5, 6.5, 10, 0.4, size=13, color=LIGHT_GRAY)

# ═══════════════════════════════════════
# Slide 7: Streamlit - Allocation
# ═══════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
add_rect(s, 0, 0, 0.08, 7.5, ACCENT)
add_text(s, "Streamlit — Capture Plan", 0.5, 0.4, 8, 0.6, size=36, bold=True, color=ACCENT)
# 3 metric cards
for i, (district, pct) in enumerate([("서초구", "31.2%"), ("영등포구", "35.8%"), ("중구", "33.0%")]):
    x = 0.5 + i * 4.2
    add_rect(s, x, 1.3, 3.8, 1.5, RGBColor(0x22, 0x2D, 0x5E))
    add_text(s, district, x + 0.3, 1.4, 3.2, 0.5, size=16, color=LIGHT_GRAY)
    add_text(s, pct, x + 0.3, 1.9, 3.2, 0.7, size=36, bold=True, color=WHITE)
add_text(s, "Actual vs Forecast Overlay", 0.5, 3.2, 6, 0.5, size=20, bold=True, color=WHITE)
add_rect(s, 0.5, 3.8, 7.5, 3, RGBColor(0x22, 0x2D, 0x5E))
add_text(s, "[3구 시계열 차트: Actual 실선 + Forecast 점선]", 1, 4.5, 6.5, 1, size=15, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(s, "Ablation MAPE", 8.5, 3.2, 4, 0.5, size=20, bold=True, color=WHITE)
add_rect(s, 8.5, 3.8, 4, 3, RGBColor(0x22, 0x2D, 0x5E))
add_text(s, "[Model A→E\nMAPE 개선 차트]", 9, 4.5, 3, 1, size=15, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# Slide 8: Streamlit - Analysis + AI
# ═══════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
add_rect(s, 0, 0, 0.08, 7.5, ACCENT)
add_text(s, "Streamlit — Move-in Signals + AI Playbook", 0.5, 0.4, 10.5, 0.6, size=34, bold=True, color=ACCENT)
# Analysis side
add_rect(s, 0.5, 1.3, 6, 5.5, RGBColor(0x22, 0x2D, 0x5E))
add_text(s, "Move-in Signals Tab", 0.7, 1.4, 5.5, 0.5, size=18, bold=True, color=WHITE)
add_multiline(s, [
    "KPI: 전입 세대 / 순이동 / 카드소비 / 자산",
    "연령구조: 20-39세, 60세+ 비중",
    "관광수요: 수요지수, 외래객 지수",
    "상권건강: 안정도, 순 점포 변동",
    "Feature Importance 차트",
    "3구 비교 테이블",
], 0.7, 2.1, 5.5, 3.5, size=14, color=ICE_BLUE)
# AI side
add_rect(s, 7, 1.3, 5.5, 5.5, RGBColor(0x22, 0x2D, 0x5E))
add_text(s, "AI Playbook Tab", 7.2, 1.4, 5, 0.5, size=18, bold=True, color=WHITE)
add_text(s, "AI_COMPLETE Structured Output:", 7.2, 2.1, 5, 0.4, size=14, bold=True, color=ACCENT)
add_multiline(s, [
    '  "recommended_district": "영등포구"',
    '  "allocation_pct": 42.0',
    '  "drivers": ["전입 증가", "카드소비 강세"]',
    '  "risk": "중구 관광 편중 변동성"',
    '  "next_action": "설치 인력과 체험 오퍼 우선 배치"',
], 7.2, 2.7, 5, 2.5, size=12, color=GREEN, bold=False)
add_text(s, "Feature Mart JSON 컨텍스트 주입\n→ Grounded Recommendation", 7.2, 5.5, 5, 0.8, size=13, color=ORANGE)

# ═══════════════════════════════════════
# Slide 9: Snowflake Expertise
# ═══════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
add_rect(s, 0, 0, 0.08, 7.5, ACCENT)
add_text(s, "Snowflake Expertise", 0.5, 0.4, 8, 0.6, size=36, bold=True, color=ACCENT)
features = [
    ("Semantic View", "DISTRICTPILOT_SV — 비즈니스 메트릭/차원/관계 정의, AI 지침 포함"),
    ("Cortex Analyst", "Semantic View 기반 정확한 SQL 생성, RBAC 연동"),
    ("Cortex Search", "정책/룰북 하이브리드 검색 (벡터 + 키워드)"),
    ("AI_COMPLETE", "Structured output (JSON schema) + CORTEX.COMPLETE fallback"),
    ("ML FORECAST", "외생변수 + evaluation + feature importance + ablation"),
    ("Dynamic Tables", "DT_FEATURE_MART + FEATURE_MART_V2 serving layer"),
    ("Tasks", "매일 06:00 갱신 + 주 1회 재학습 + 자동 재시도"),
    ("Streamlit in SiS", "Owner's rights, 데이터 외부 미유출, RBAC 접근 관리"),
]
for i, (name, desc) in enumerate(features):
    y = 1.2 + i * 0.72
    add_rect(s, 0.5, y, 12, 0.62, RGBColor(0x22, 0x2D, 0x5E) if i % 2 == 0 else RGBColor(0x1A, 0x23, 0x50))
    add_text(s, name, 0.7, y + 0.05, 2.5, 0.5, size=14, bold=True, color=ACCENT)
    add_text(s, desc, 3.3, y + 0.05, 8.8, 0.5, size=13, color=ICE_BLUE)

# ═══════════════════════════════════════
# Slide 10: Operations
# ═══════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
add_rect(s, 0, 0, 0.08, 7.5, ACCENT)
add_text(s, "Operations / Trust", 0.5, 0.4, 6, 0.6, size=36, bold=True, color=ACCENT)
add_table_slide(s, ["Component", "Type", "Status"], [
    ["DT_FEATURE_MART", "Dynamic Table", "1h target lag"],
    ["FEATURE_MART_V2", "Serving Table", "App input"],
    ["TASK_REFRESH_PIPELINE", "Task", "Daily 06:00 KST"],
    ["TASK_RETRAIN_FORECAST", "Task", "Weekly Mon 07:00"],
    ["V_APP_HEALTH", "View", "Real-time monitoring"],
    ["DISTRICTPILOT_FORECAST_V2", "ML Model", "Ablation E (production)"],
], 0.5, 1.3, 7.5)
add_rect(s, 8.5, 1.3, 4, 5.2, RGBColor(0x22, 0x2D, 0x5E))
add_text(s, "Security Model", 8.7, 1.5, 3.6, 0.4, size=16, bold=True, color=ACCENT)
add_multiline(s, [
    "Streamlit: Owner's rights",
    "Cortex: RBAC 연동",
    "Query tag: JSON 추적",
    "데이터: Marketplace 내부 복제",
    "",
    "Governance",
    "스폰서 데이터 미공유/미배포",
    "외부 데이터: 공공누리 1유형",
    "라이선스 링크 코드 내 명시",
], 8.7, 2.1, 3.6, 4, size=13, color=ICE_BLUE)

# ═══════════════════════════════════════
# Slide 11: One Engine, One Operating Loop
# ═══════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
add_rect(s, 0, 0, 0.08, 7.5, ACCENT)
add_text(s, "One Engine, One Operating Loop", 0.5, 0.4, 10.5, 0.6, size=34, bold=True, color=ACCENT)
add_text(s, "같은 엔진이 획득 전략과 운영 실행을 동시에 연결", 0.5, 1.0, 10, 0.5, size=18, color=ICE_BLUE)
# Acquisition
add_rect(s, 0.5, 1.8, 5.8, 4.8, RGBColor(0x22, 0x2D, 0x5E))
add_rect(s, 0.5, 1.8, 5.8, 0.06, ACCENT)
add_text(s, "획득 전략", 0.7, 2.0, 5.4, 0.5, size=22, bold=True, color=ACCENT)
add_multiline(s, [
    "전입 직후 타깃 캠페인 우선순위",
    "구별 집행 강도와 채널 믹스",
    "체험 오퍼/번들 제안",
    "B2B/B2C 권역 공략 순서",
    "월간 시나리오 비교",
], 0.8, 2.7, 5.2, 3.5, size=16, color=WHITE)
# Operations
add_rect(s, 6.8, 1.8, 5.8, 4.8, RGBColor(0x22, 0x2D, 0x5E))
add_rect(s, 6.8, 1.8, 5.8, 0.06, ORANGE)
add_text(s, "운영 실행", 7.0, 2.0, 5.4, 0.5, size=22, bold=True, color=ORANGE)
add_multiline(s, [
    "설치 가능 지역/시간대 점검",
    "현장 인력·재고 선배치",
    "오피스/관광 권역 운영 제약 반영",
    "CS 에스컬레이션과 고위험 구역 관리",
    "실험 예산과 후속 조정",
], 7.1, 2.7, 5.2, 3.5, size=16, color=WHITE)

# ═══════════════════════════════════════
# Slide 12: Cost
# ═══════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
add_rect(s, 0, 0, 0.08, 7.5, ACCENT)
add_text(s, "Cost", 0.5, 0.4, 5, 0.6, size=36, bold=True, color=ACCENT)
add_text(s, "~$80/month", 0.5, 1.2, 6, 1, size=48, bold=True, color=WHITE)
add_table_slide(s, ["Component", "Spec", "Est. Cost"], [
    ["Compute WH", "X-Small, auto-suspend", "~$30"],
    ["Cortex LLM", "mistral-large2, ~50 calls/day", "~$25"],
    ["Streamlit hosting", "SiS included", "~$5"],
    ["Dynamic Tables", "1h refresh, 2 DTs", "~$10"],
    ["Tasks", "Daily + Weekly", "~$5"],
    ["Storage", "~1GB total", "~$5"],
], 0.5, 2.5, 7)

# ═══════════════════════════════════════
# Slide 13: Next Steps
# ═══════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
add_rect(s, 0, 0, 0.08, 7.5, ACCENT)
add_text(s, "Next Steps", 0.5, 0.4, 5, 0.6, size=36, bold=True, color=ACCENT)
steps = [
    ("Phase 1", "AJD 선택 통합", "상품/채널 수준 추천까지 확장 가능한 옵션 레이어"),
    ("Phase 2", "Cortex Agent", "Analyst + Search + Custom tool 오케스트레이션"),
    ("Phase 3", "MCP Server", "Managed MCP로 외부 시스템 안전 연결"),
    ("Phase 4", "Native App", "Marketplace listing으로 앱 배포"),
    ("Phase 5", "25개 구 확장", "서울 전체 → 수도권 → 전국"),
]
for i, (phase, title, desc) in enumerate(steps):
    y = 1.3 + i * 1.1
    add_rect(s, 0.5, y, 12, 0.9, RGBColor(0x22, 0x2D, 0x5E))
    add_rect(s, 0.5, y, 0.06, 0.9, ACCENT if i < 2 else ORANGE if i < 4 else GREEN)
    add_text(s, phase, 0.8, y + 0.1, 1.5, 0.4, size=14, bold=True, color=ACCENT)
    add_text(s, title, 2.5, y + 0.05, 3, 0.4, size=17, bold=True, color=WHITE)
    add_text(s, desc, 2.5, y + 0.45, 9.5, 0.4, size=13, color=ICE_BLUE)

# ═══════════════════════════════════════
# Slide 14: Tech Stack Summary
# ═══════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
add_rect(s, 0, 0, 0.08, 7.5, ACCENT)
add_text(s, "Tech Stack", 0.5, 0.4, 5, 0.6, size=36, bold=True, color=ACCENT)
add_table_slide(s, ["Object", "Type", "Role"], [
    ["FEATURE_MART_V2", "Table", "확장 Feature Mart (50+ features)"],
    ["DT_FEATURE_MART", "Dynamic Table", "자동 갱신 (1h lag)"],
    ["DISTRICTPILOT_FORECAST_V2", "ML Model", "Ablation E (production)"],
    ["DISTRICTPILOT_SV", "Semantic View", "Cortex Analyst용 메트릭 정의"],
    ["DISTRICTPILOT_SEARCH_SVC", "Cortex Search", "정책/룰북 검색"],
    ["ABLATION_RESULTS", "Table/View", "5 모델 MAPE/SMAPE/MAE"],
    ["FORECAST_RESULTS", "Table", "예측 결과 serving"],
    ["TASK_REFRESH_PIPELINE", "Task", "매일 06:00 KST"],
    ["TASK_RETRAIN_FORECAST", "Task", "매주 월 07:00 KST"],
    ["V_APP_HEALTH", "View", "운영 상태 모니터링"],
    ["POLICY_DOCUMENTS", "Table", "AI 근거 정책 문서"],
], 0.5, 1.1, 12)

# ═══════════════════════════════════════
# Slide 15: Evidence Chain
# ═══════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
add_rect(s, 0, 0, 0.08, 7.5, ACCENT)
add_text(s, "Evidence Chain", 0.5, 0.4, 8, 0.6, size=36, bold=True, color=ACCENT)
add_text(s, "의사결정의 6단계 증거 — 모두 Snowflake-native", 0.5, 1.0, 10, 0.5, size=18, color=ICE_BLUE)

chain_steps = [
    ("1", "Forecast", "DISTRICTPILOT_FORECAST_V2", "3구 x 다월 예측, 평가 지표와 함께 검증"),
    ("2", "Feature Importance", "EXPLAIN_FEATURE_IMPORTANCE()", "Top-5 기여도 + ablation 개선 증명"),
    ("3", "Semantic View", "DISTRICTPILOT_SV", "VQR 10개, SQL규칙 13개, synonym 100+"),
    ("4", "Search Grounding", "CORTEX.SEARCH()", "외부 컨텍스트 근거, hallucination 방지"),
    ("5", "Structured Output", "AI_COMPLETE() + CORTEX.COMPLETE()", "JSON action card + confidence score"),
    ("6", "Refresh State", "V_APP_HEALTH", "LAG_SEC, 갱신 이력, task 상태, query_tag 감사"),
]
for i, (num, title, obj, desc) in enumerate(chain_steps):
    y = 1.7 + i * 0.85
    add_rect(s, 0.5, y, 12, 0.75, RGBColor(0x22, 0x2D, 0x5E))
    add_rect(s, 0.5, y, 0.06, 0.75, ACCENT)
    # Circle number
    add_text(s, num, 0.7, y + 0.08, 0.5, 0.5, size=20, bold=True, color=ACCENT)
    add_text(s, title, 1.3, y + 0.08, 2.5, 0.35, size=15, bold=True, color=WHITE)
    add_text(s, obj, 4.0, y + 0.08, 3.5, 0.35, size=12, color=GREEN)
    add_text(s, desc, 7.5, y + 0.08, 5, 0.6, size=12, color=ICE_BLUE)

add_text(s, "모든 단계가 Snowflake 안에서 완결 — 데이터가 밖으로 나가지 않습니다", 0.5, 6.8, 12, 0.4, size=14, bold=True, color=ORANGE)

# ═══════════════════════════════════════
# Slide 16: Q&A
# ═══════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, DARK_NAVY)
add_rect(s, 0, 0, 13.333, 0.08, ACCENT)
add_text(s, "감사합니다", 0.5, 2.0, 12.3, 1.5, size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "DistrictPilot AI — Snowflake-Native Move-in Demand Orchestration Engine", 0.5, 3.5, 12.3, 0.8, size=20, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s, "Doeon Kim  |  github.com/KIM3310", 0.5, 4.8, 12.3, 0.5, size=16, color=ICE_BLUE, align=PP_ALIGN.CENTER)
add_text(s, "Q & A", 0.5, 5.8, 12.3, 0.8, size=36, bold=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Save
outpath = "/Users/pizza/districtpilot-ai/deliverables/DistrictPilot_AI_Hackathon_Final.pptx"
prs.save(outpath)
print(f"Saved to {outpath}")
print(f"Slides: {len(prs.slides)}")
